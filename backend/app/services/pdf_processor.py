"""
PDF Processing Service (Docling)
==================================
Handles the extraction of structured text from uploaded PDF documents.

WHY DOCLING (not PyPDF2, pdfplumber, or pdfminer):
University exam papers are structurally complex:
  - Multi-column layouts (common in Indian university papers)
  - Tables with mark allocations, sub-question grids
  - Mixed fonts, bold headings, nested numbering (4.a.i, 4.a.ii)

Generic PDF extractors read text LEFT-TO-RIGHT across the full page width.
On a two-column paper, this means the first line of column A gets concatenated
with the first line of column B — destroying all meaning.

Docling uses IBM's layout analysis AI model (DocLayNet) to:
  1. Detect page regions (text blocks, tables, headers, footers)
  2. Determine reading order (column A fully, then column B)
  3. Reconstruct tables as proper Markdown tables
  4. Preserve heading hierarchy for chunk boundary detection

THE RESULT: Clean, correctly-ordered Markdown that maintains the
relational structure between questions, sub-questions, and mark allocations.
"""

import os
import logging
from pathlib import Path
from typing import Optional
from app.utils.text_utils import clean_extracted_text

logger = logging.getLogger(__name__)


class PDFProcessor:
    """
    Extracts structured content from PDF files using Docling.
    
    LIFECYCLE:
    1. __init__: Creates the Docling DocumentConverter (loads the AI model ~once)
    2. process_file: Takes a file path, returns clean Markdown + metadata
    3. The caller (upload router) saves the Markdown to the database
    """

    def __init__(self):
        """
        Initialize the Docling converter.
        
        NOTE: The first call loads the layout model into memory (~500MB).
        Subsequent calls reuse the loaded model. On free-tier hosting,
        this initial load takes ~10-15 seconds. After that, each PDF
        takes ~2-8 seconds depending on page count.
        """
        try:
            from docling.document_converter import DocumentConverter
            self.converter = DocumentConverter()
            self._available = True
            logger.info("Docling PDF processor initialized successfully")
        except ImportError:
            # Docling might not be installed in dev/test environments
            self.converter = None
            self._available = False
            logger.warning(
                "Docling is not installed. PDF processing will use fallback extraction. "
                "Install with: pip install docling"
            )

    @property
    def is_available(self) -> bool:
        """Check if Docling is loaded and ready."""
        return self._available

    async def process_file(self, file_path: str) -> dict:
        """
        Extract structured text from a PDF file.
        
        Args:
            file_path: Absolute path to the PDF file on disk
        
        Returns:
            {
                "markdown": str,      # Full document as clean Markdown
                "tables": list[str],  # Each table as a Markdown table string
                "page_count": int,    # Number of pages processed
                "headings": list[str] # Detected section headings
            }
        
        HOW THE PIPELINE FLOWS:
        file_path → Docling converter → DoclingDocument → Markdown export
                                                        → Table extraction
                                                        → Heading detection
        """
        if not self._available:
            return await self._fallback_extract(file_path)

        try:
            logger.info(f"Starting Docling extraction for: {Path(file_path).name}")
            
            # Docling's convert() handles all the heavy lifting:
            # Layout detection → Reading order → Table recognition → Structure mapping
            result = self.converter.convert(file_path)
            doc = result.document

            # Export the full document as Markdown
            # Docling preserves: headings, lists, tables, bold/italic, code blocks
            raw_markdown = doc.export_to_markdown()
            
            # Clean up common PDF artifacts (page numbers, excessive whitespace)
            markdown = clean_extracted_text(raw_markdown)

            # Extract tables separately for structured analysis
            # Each table becomes a proper Markdown table: | Col1 | Col2 | ...
            tables = []
            if hasattr(doc, 'tables') and doc.tables:
                for table in doc.tables:
                    try:
                        table_md = table.export_to_markdown()
                        if table_md.strip():
                            tables.append(table_md)
                    except Exception as e:
                        logger.warning(f"Failed to export table: {e}")

            # Extract headings for chunk boundary detection
            headings = self._extract_headings(markdown)

            # Count pages
            page_count = 0
            if hasattr(result, 'pages') and result.pages:
                page_count = len(result.pages)
            elif hasattr(doc, 'pages') and doc.pages:
                page_count = len(doc.pages)

            # Fallback for PPTX/PDF if Docling doesn't extract page counts properly
            if page_count == 0:
                file_path_str = str(file_path).lower()
                if file_path_str.endswith('.pptx') or file_path_str.endswith('.ppt'):
                    try:
                        from pptx import Presentation
                        prs = Presentation(file_path)
                        page_count = len(prs.slides)
                    except Exception:
                        pass
                elif file_path_str.endswith('.pdf'):
                    try:
                        import PyPDF2
                        with open(file_path, 'rb') as f:
                            reader = PyPDF2.PdfReader(f)
                            page_count = len(reader.pages)
                    except Exception:
                        pass

            logger.info(
                f"Extraction complete: {page_count} pages, "
                f"{len(tables)} tables, {len(headings)} headings"
            )

            return {
                "markdown": markdown,
                "tables": tables,
                "page_count": page_count,
                "headings": headings
            }

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            logger.error(f"Docling extraction failed for {file_path}: {e}\n{error_details}")
            try:
                with open(f"docling_error.log", "w") as f:
                    f.write(error_details)
            except Exception:
                pass
            # Fall back to basic extraction rather than failing completely
            return await self._fallback_extract(file_path)

    def _extract_headings(self, markdown: str) -> list[str]:
        """
        Extract Markdown headings from the document.
        These are used to identify natural chunk boundaries.
        
        Heading patterns: # Title, ## Section, ### Subsection
        """
        import re
        headings = []
        for line in markdown.split('\n'):
            line = line.strip()
            if line.startswith('#'):
                # Remove the # symbols and clean up
                heading = re.sub(r'^#+\s*', '', line).strip()
                if heading and len(heading) > 2:
                    headings.append(heading)
        return headings

    async def _fallback_extract(self, file_path: str) -> dict:
        """
        Basic text extraction when Docling is unavailable.
        
        Uses PyPDF2 as a minimal fallback — this won't handle tables or
        multi-column layouts correctly, but at least extracts raw text
        so the pipeline doesn't break completely during development.
        """
        logger.warning(f"Using fallback PDF extraction for: {file_path}")
        
        try:
            file_path_str = str(file_path).lower()
            text_parts = []
            page_count = 0
            
            if file_path_str.endswith('.pdf'):
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    page_count = len(reader.pages)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                            
            elif file_path_str.endswith('.pptx') or file_path_str.endswith('.ppt'):
                from pptx import Presentation
                prs = Presentation(file_path)
                page_count = len(prs.slides)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                            
            elif file_path_str.endswith('.docx') or file_path_str.endswith('.doc'):
                import docx
                doc = docx.Document(file_path)
                page_count = 1 # DOCX doesn't have a fixed page count easily
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                        
            else:
                logger.error(f"Fallback extraction unsupported format for {file_path}")
                return {"markdown": "", "tables": [], "page_count": 0, "headings": []}
            
            markdown = clean_extracted_text('\n\n'.join(text_parts))
            
            return {
                "markdown": markdown,
                "tables": [],
                "page_count": page_count,
                "headings": self._extract_headings(markdown)
            }
        except Exception as e:
            logger.error(f"Fallback extraction failed: {e}")
            return {
                "markdown": "",
                "tables": [],
                "page_count": 0,
                "headings": []
            }


# Singleton instance — reuse across requests to avoid reloading the AI model
_processor: Optional[PDFProcessor] = None


def get_pdf_processor() -> PDFProcessor:
    """Get or create the global PDF processor instance."""
    global _processor
    if _processor is None:
        _processor = PDFProcessor()
    return _processor
